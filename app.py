from flask import Flask, request, jsonify, send_file
from flask_cors import CORS
import google.genai as genai
import json, io, os, re
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.lib.colors import HexColor
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
from docx import Document
from docx.shared import Pt, RGBColor, Cm
from docx.enum.text import WD_ALIGN_PARAGRAPH
import openpyxl
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN

app = Flask(__name__)
CORS(app, origins='*', methods=['GET', 'POST', 'OPTIONS'], allow_headers=['Content-Type'])

GEMINI_KEY = os.environ.get('GOOGLE_API_KEY') or os.environ.get('GEMINI_API_KEY', '')

SYSTEM_PROMPT = """Você é Maia, agente especialista em criação de materiais profissionais da Mentoria Âncora, criada e treinada por Jhenifer.

Sua função é criar materiais prontos para download: PDF, Word, Excel e PowerPoint.

QUEM VOCÊ ATENDE
Empresárias de todos os segmentos: comércio local, loja física, infoproduto, serviço, mentoria, produto físico.

ANTES DE CRIAR
Faça no máximo 2 perguntas essenciais antes de criar. Se o pedido já tem informação suficiente, crie direto.

O QUE VOCÊ CRIA
Planilhas financeiras, apresentações de vendas, contratos, planejamentos, cronogramas, materiais didáticos.

QUANDO FOR CRIAR UM ARQUIVO, responda SEMPRE neste formato JSON exato, sem nenhum texto fora dele:

Para PDF e Word:
{"tipo":"pdf","nome_arquivo":"nome-sem-extensao","mensagem":"mensagem curta","conteudo":{"titulo":"Título","secoes":[{"titulo":"Título da Seção","texto":"Conteúdo detalhado da seção"}]}}

Para Excel:
{"tipo":"excel","nome_arquivo":"nome-sem-extensao","mensagem":"mensagem curta","conteudo":{"planilhas":[{"nome":"Aba","cabecalhos":["Col1","Col2"],"linhas":[["dado1","dado2"]]}]}}

Para PowerPoint:
{"tipo":"pptx","nome_arquivo":"nome-sem-extensao","mensagem":"mensagem curta","conteudo":{"titulo_apresentacao":"Título","slides":[{"titulo":"Slide","pontos":["ponto 1","ponto 2"]}]}}

REGRAS ABSOLUTAS
Nunca use travessão. Nunca entregue conteúdo genérico. Nunca use linguagem formal ou distante. Sempre crie materiais completos. Quando for saudação ou pergunta simples, responda em texto puro sem JSON."""

def get_client():
    return genai.Client(api_key=GEMINI_KEY)

def criar_pdf(conteudo, nome):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4,
        rightMargin=2*cm, leftMargin=2*cm,
        topMargin=2.5*cm, bottomMargin=2.5*cm)
    preto = HexColor('#0d0d0d')
    cinza_escuro = HexColor('#1a1a1a')
    cinza_medio = HexColor('#444444')
    titulo_style = ParagraphStyle('titulo', fontName='Helvetica-Bold', fontSize=22, textColor=preto, spaceAfter=8, leading=28)
    secao_style = ParagraphStyle('secao', fontName='Helvetica-Bold', fontSize=12, textColor=cinza_escuro, spaceBefore=16, spaceAfter=6, leading=16)
    corpo_style = ParagraphStyle('corpo', fontName='Helvetica', fontSize=10, textColor=cinza_medio, spaceAfter=6, leading=16)
    elementos = []
    elementos.append(Paragraph(conteudo.get('titulo', nome), titulo_style))
    elementos.append(HRFlowable(width='100%', thickness=1.5, color=preto, spaceAfter=16))
    for secao in conteudo.get('secoes', []):
        elementos.append(Paragraph(secao.get('titulo', ''), secao_style))
        for linha in secao.get('texto', '').split('\n'):
            if linha.strip():
                elementos.append(Paragraph(linha.strip(), corpo_style))
        elementos.append(Spacer(1, 6))
    doc.build(elementos)
    buffer.seek(0)
    return buffer

def criar_word(conteudo, nome):
    doc = Document()
    for section in doc.sections:
        section.top_margin = Cm(2.5)
        section.bottom_margin = Cm(2.5)
        section.left_margin = Cm(2.5)
        section.right_margin = Cm(2.5)
    titulo = doc.add_heading(conteudo.get('titulo', nome), 0)
    titulo.alignment = WD_ALIGN_PARAGRAPH.LEFT
    for run in titulo.runs:
        run.font.color.rgb = RGBColor(13, 13, 13)
        run.font.size = Pt(22)
    doc.add_paragraph()
    for secao in conteudo.get('secoes', []):
        h = doc.add_heading(secao.get('titulo', ''), 1)
        for run in h.runs:
            run.font.color.rgb = RGBColor(26, 26, 26)
            run.font.size = Pt(12)
        for linha in secao.get('texto', '').split('\n'):
            if linha.strip():
                p = doc.add_paragraph(linha.strip())
                for run in p.runs:
                    run.font.color.rgb = RGBColor(68, 68, 68)
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def criar_excel(conteudo, nome):
    wb = openpyxl.Workbook()
    wb.remove(wb.active)
    header_fill = PatternFill(start_color='0D0D0D', end_color='0D0D0D', fill_type='solid')
    alt_fill = PatternFill(start_color='F5F5F5', end_color='F5F5F5', fill_type='solid')
    border = Border(bottom=Side(style='thin', color='DDDDDD'))
    for planilha in conteudo.get('planilhas', []):
        ws = wb.create_sheet(title=planilha.get('nome', 'Planilha'))
        for col, cab in enumerate(planilha.get('cabecalhos', []), 1):
            cell = ws.cell(row=1, column=col, value=cab)
            cell.font = Font(bold=True, color='FFFFFF', size=11)
            cell.fill = header_fill
            cell.alignment = Alignment(horizontal='center', vertical='center')
            ws.column_dimensions[cell.column_letter].width = max(18, len(str(cab)) + 4)
        for row_idx, linha in enumerate(planilha.get('linhas', []), 2):
            for col_idx, valor in enumerate(linha, 1):
                cell = ws.cell(row=row_idx, column=col_idx, value=valor)
                cell.alignment = Alignment(horizontal='left', vertical='center')
                cell.border = border
                if row_idx % 2 == 0:
                    cell.fill = alt_fill
        ws.row_dimensions[1].height = 28
    buffer = io.BytesIO()
    wb.save(buffer)
    buffer.seek(0)
    return buffer

def criar_pptx(conteudo, nome):
    prs = Presentation()
    prs.slide_width = Emu(9144000)
    prs.slide_height = Emu(5143500)
    titulo_apres = conteudo.get('titulo_apresentacao', nome)
    for i, slide_data in enumerate(conteudo.get('slides', [])):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = PptRGB(13, 13, 13) if i == 0 else PptRGB(255, 255, 255)
        if i == 0:
            tb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(1.2))
            p = tb.text_frame.paragraphs[0]
            p.text = titulo_apres
            p.font.size = PptPt(36)
            p.font.bold = True
            p.font.color.rgb = PptRGB(240, 236, 228)
            p.alignment = PP_ALIGN.CENTER
        else:
            tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
            p = tb.text_frame.paragraphs[0]
            p.text = slide_data.get('titulo', '')
            p.font.size = PptPt(22)
            p.font.bold = True
            p.font.color.rgb = PptRGB(13, 13, 13)
            cb = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(3.5))
            cb.text_frame.word_wrap = True
            for j, ponto in enumerate(slide_data.get('pontos', [])):
                p2 = cb.text_frame.paragraphs[0] if j == 0 else cb.text_frame.add_paragraph()
                p2.text = f"  {ponto}"
                p2.font.size = PptPt(14)
                p2.font.color.rgb = PptRGB(50, 50, 50)
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

@app.route('/chat', methods=['POST', 'OPTIONS'])
def chat():
    if request.method == 'OPTIONS':
        return '', 200
    data = request.json
    messages = data.get('messages', [])
    history = []
    for m in messages[:-1]:
        role = 'model' if m['role'] == 'assistant' else m['role']
        history.append({'role': role, 'parts': [{'text': m['content']}]})
    last = messages[-1]['content'] if messages else ''
    response = get_client().models.generate_content(
        model='gemini-2.5-flash',
        contents=history + [{'role': 'user', 'parts': [{'text': last}]}],
        config={'system_instruction': SYSTEM_PROMPT, 'temperature': 0.7, 'max_output_tokens': 4000}
    )
    reply = response.text.strip()
    try:
        clean = re.sub(r'```json|```', '', reply).strip()
        parsed = json.loads(clean)
        if parsed.get('tipo') and parsed.get('conteudo'):
            return jsonify({'type': 'file_ready', 'tipo': parsed['tipo'],
                          'nome': parsed['nome_arquivo'], 'mensagem': parsed['mensagem'],
                          'conteudo': parsed['conteudo']})
    except:
        pass
    return jsonify({'type': 'text', 'message': reply})

@app.route('/gerar', methods=['POST', 'OPTIONS'])
def gerar():
    if request.method == 'OPTIONS':
        return '', 200
    data = request.json
    tipo = data.get('tipo')
    nome = data.get('nome', 'arquivo')
    conteudo = data.get('conteudo', {})
    extensoes = {'pdf': 'pdf', 'word': 'docx', 'excel': 'xlsx', 'pptx': 'pptx'}
    mime_types = {'pdf': 'application/pdf',
                  'word': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                  'excel': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
                  'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'}
    if tipo == 'pdf': buffer = criar_pdf(conteudo, nome)
    elif tipo == 'word': buffer = criar_word(conteudo, nome)
    elif tipo == 'excel': buffer = criar_excel(conteudo, nome)
    elif tipo == 'pptx': buffer = criar_pptx(conteudo, nome)
    else: return jsonify({'erro': 'Tipo inválido'}), 400
    return send_file(buffer, mimetype=mime_types[tipo],
                    as_attachment=True, download_name=f'{nome}.{extensoes[tipo]}')

@app.route('/health', methods=['GET'])
def health():
    return jsonify({'status': 'ok'})

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=int(os.environ.get('PORT', 5000)))
